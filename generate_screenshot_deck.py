from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "ai_agent_skills_screenshots.pptx"
QA_OUT = ROOT / "qa_preview_ai_agent_skills_screenshots.png"

SLIDE_W = 13.333
SLIDE_H = 7.5

IMAGES = [
    {
        "path": ROOT / "Screenshot 2026-05-05 at 8.06.46 PM.png",
        "title": "Models Are Strong. Context Still Matters.",
        "caption": "Opening frame: model quality improves, but the harness and context shape the output.",
    },
    {
        "path": ROOT / "Screenshot 2026-05-05 at 8.07.22 PM.png",
        "title": "Context Stack for Agent Work",
        "caption": "Global instructions, agent files, prompts, tools, user request, and skills all contribute context.",
    },
    {
        "path": ROOT / "unnamed (5).png",
        "title": "Anatomy of AI Agent Skills",
        "caption": "Infographic view: skill.md, scripts, references, assets, and progressive disclosure.",
    },
    {
        "path": ROOT / "unnamed (4).png",
        "title": "Procedural Memory and Skill Loading",
        "caption": "Second infographic source: skills as portable procedural memory, distinct from MCP, RAG, and fine-tuning.",
    },
]


def set_text_style(run, size, color="FFFFFF", bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.name = "Aptos"


def rgb(hex_color):
    from pptx.dml.color import RGBColor

    hex_color = hex_color.strip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def add_textbox(slide, text, x, y, w, h, size, color, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.font.name = "Aptos"
    return box


def add_image_contain(slide, path, x, y, w, h, border_color="D8DEE9"):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    draw_w = iw * scale
    draw_h = ih * scale
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2

    pad = 0.03
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(draw_x - pad),
        Inches(draw_y - pad),
        Inches(draw_w + pad * 2),
        Inches(draw_h + pad * 2),
    ).fill.solid()
    frame = slide.shapes[-1]
    frame.fill.fore_color.rgb = rgb("FFFFFF")
    frame.line.color.rgb = rgb(border_color)
    frame.line.width = Pt(0.75)

    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))
    return draw_x, draw_y, draw_w, draw_h


def add_deck_header(slide, title, caption):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("F7F8FA")
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.78)
    ).fill.solid()
    header = slide.shapes[-1]
    header.fill.fore_color.rgb = rgb("18202A")
    header.line.fill.background()
    add_textbox(slide, title, 0.45, 0.18, 8.4, 0.4, 20, "FFFFFF", True)
    add_textbox(slide, caption, 8.1, 0.22, 4.75, 0.34, 9.5, "B8C1CC", False)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for item in IMAGES:
        slide = prs.slides.add_slide(blank)
        add_deck_header(slide, item["title"], item["caption"])
        add_image_contain(slide, item["path"], 0.35, 1.02, 12.63, 6.08)
        add_textbox(slide, item["path"].name, 0.45, 7.14, 9.5, 0.18, 8, "66717F")

    prs.save(OUT)


def build_qa_preview():
    scale = 120
    canvas_w = int(SLIDE_W * scale)
    canvas_h = int(SLIDE_H * scale)
    cols = 2
    rows = 2
    pad = 28
    preview = Image.new("RGB", (cols * canvas_w + (cols + 1) * pad, rows * canvas_h + (rows + 1) * pad), "white")
    draw = ImageDraw.Draw(preview)
    try:
        font_title = ImageFont.truetype("Arial.ttf", 28)
        font_small = ImageFont.truetype("Arial.ttf", 16)
    except OSError:
        font_title = ImageFont.load_default()
        font_small = ImageFont.load_default()

    slides = [(item["title"], item["path"]) for item in IMAGES]
    for idx, (title, path) in enumerate(slides):
        col = idx % cols
        row = idx // cols
        ox = pad + col * (canvas_w + pad)
        oy = pad + row * (canvas_h + pad)
        draw.rectangle([ox, oy, ox + canvas_w, oy + canvas_h], fill="#F7F8FA", outline="#D8DEE9", width=2)
        draw.rectangle([ox, oy, ox + canvas_w, oy + int(0.78 * scale)], fill="#18202A")
        draw.text((ox + 54, oy + 22), title, fill="white", font=font_small)
        with Image.open(path) as img:
            img = img.convert("RGB")
            max_w = int(12.63 * scale)
            max_h = int(6.08 * scale)
            img.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
            ix = ox + int(0.35 * scale) + (max_w - img.width) // 2
            iy = oy + int(1.02 * scale) + (max_h - img.height) // 2
            draw.rectangle([ix - 4, iy - 4, ix + img.width + 4, iy + img.height + 4], fill="white", outline="#D8DEE9")
            preview.paste(img, (ix, iy))

    preview.save(QA_OUT)


if __name__ == "__main__":
    build_deck()
    build_qa_preview()
    print(OUT)
    print(QA_OUT)
